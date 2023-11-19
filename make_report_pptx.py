# Libraries
# pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# pandas
import pandas as pd
import numpy as np

# dataset sample
import sample_dataset as sd

# parameter
orange = (214, 87, 20)
blue = (0,0,128)
sky = (0,153,255)

# class
class MakePPT:

    def __init__(self, savefilename:str):

        # init
        self.savefilename = savefilename
        # new presenntation
        self.prs = Presentation()
        # wide screen
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    ##########################################
    # 1st slide, analysis data summary
    ##########################################
    def data_summary(self,
                     dataset_summary:pd.DataFrame,
                     ml_learning_dataset_summary:pd.DataFrame):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # ------------------------
        # title
        # ------------------------
        title = slide.shapes.title
        title.text = "Analysis data summary"
        # title text format and position
        title = self.format_title(title=title)

        # ------------------------
        # dataset summary
        # ------------------------
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(Inches(0.5),
                                            Inches(1),
                                            Inches(5),
                                            Inches(0.3))
        # add text
        textbox1.text_frame.paragraphs[0].text = "Dataset summary"
        textbox1.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position
        # add dataframe
        self.add_matrix_1(slide=slide,
                        df=dataset_summary,
                        left=0.5,
                        top=1.5,
                        width=5,
                        height=10,
                        column_color=orange,
                        merge_columns=[(0,"Item")])

        # ------------------------
        # ml learning dataset summary
        # ------------------------
        # add text box, left -> top -> width -> height
        textbox2 = slide.shapes.add_textbox(Inches(6.5),
                                            Inches(1),
                                            Inches(5),
                                            Inches(0.3))
        # add text
        textbox2.text_frame.paragraphs[0].text = "ML Learning dataset summary"
        textbox2.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox2.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox2.text_frame.paragraphs[0].font.bold = True # bold
        textbox2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position
        # add dataframe
        self.add_matrix_1(slide=slide,
                        df=ml_learning_dataset_summary,
                        left=6.5,
                        top=1.5,
                        width=5,
                        height=10,
                        column_color=orange,
                        merge_columns=[(0,"Item")])

    ##########################################
    # 2nd slide, Yield diagram
    ##########################################
    def yield_diagram(self,
                     diagram:pd.DataFrame):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # ------------------------
        # title
        # ------------------------
        title = slide.shapes.title
        title.text = "Yield diagram"
        # title text format and position
        title = self.format_title(title=title)

        # ------------------------
        # diagram
        # ------------------------
        # text box 1
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(Inches(0.4),
                                            Inches(0.5),
                                            Inches(4),
                                            Inches(0.3))
        # add text
        textbox1.text_frame.paragraphs[0].text = "Failure rate analysis (Pareto analysis)"
        textbox1.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # text box 2
        # add text box, left -> top -> width -> height
        textbox2 = slide.shapes.add_textbox(Inches(4.25),
                                            Inches(0.5),
                                            Inches(4),
                                            Inches(0.3))
        # add text
        textbox2.text_frame.paragraphs[0].text = "Auto Condition Analysis"
        textbox2.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox2.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox2.text_frame.paragraphs[0].font.bold = True # bold
        textbox2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # text box 3
        # add text box, left -> top -> width -> height
        textbox3 = slide.shapes.add_textbox(Inches(9.25),
                                            Inches(0.5),
                                            Inches(4),
                                            Inches(0.3))
        # add text
        textbox3.text_frame.paragraphs[0].text = "Auto Factor Analysis (Prediction by ML)"
        textbox3.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox3.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox3.text_frame.paragraphs[0].font.bold = True # bold
        textbox3.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # diagram
        # add dataframe
        self.add_matrix_2(slide=slide,
                        df=diagram,
                        left=0.4,
                        top=0.8,
                        width=12,
                        height=15,
                        merge_columns=[(0,"Total"), (1,"Y*"), (2, "BIN")])

    ##########################################
    # 3rd slide, Test condition analysis result
    ##########################################
    def test_data_condition(self,
                            test_name:str,
                            stats_matrix:pd.DataFrame,
                            stats_histgram:str,
                            condition_analysis:pd.DataFrame,
                            map_images_df:pd.DataFrame,
                            graph_images_df:pd.DataFrame):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # ------------------------
        # title
        # ------------------------
        title = slide.shapes.title
        title.text = f"Test : {test_name}, Condition analysis result"
        # title text format and position
        title = self.format_title(title=title)

        # --------------------------
        # sub title by text box
        # --------------------------
        self.add_subtitle(slide=slide,
                          subtitle_sentence="Please input comment.")

        # --------------------------
        # Statistics value
        # --------------------------
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(Inches(0.5),
                                            Inches(1),
                                            Inches(3),
                                            Inches(0.3))
        # add text
        textbox1.text_frame.paragraphs[0].text = "Statistics value"
        textbox1.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # matrix
        # add dataframe
        self.add_matrix_1(slide=slide,
                        df=stats_matrix,
                        left=0.5,
                        top=1.5,
                        width=3,
                        height=5,
                        column_color=orange)

        # histgram image
        # add slide
        slide.shapes.add_picture(stats_histgram,
                                 Inches(0.5),
                                 Inches(4.2),
                                 Inches(3),
                                 Inches(3)
                                 )

        # --------------------------
        # Auto condition analysis
        # --------------------------
        # add text box, left -> top -> width -> height
        textbox2 = slide.shapes.add_textbox(Inches(3.6),
                                            Inches(1),
                                            Inches(3),
                                            Inches(0.3))
        # add text
        textbox2.text_frame.paragraphs[0].text = "Auto Condition Analysis Result"
        textbox2.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox2.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox2.text_frame.paragraphs[0].font.bold = True # bold
        textbox2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # matrix
        # add dataframe
        self.add_matrix_1(slide=slide,
                        df=condition_analysis,
                        left=3.6,
                        top=1.5,
                        width=3,
                        height=5,
                        column_color=blue,
                        merge_columns=None)

        # map images
        # image start position
        left = 3.6
        top = 4.2
        width = 0.7
        height = 0.7
        for idx, row in map_images_df.iterrows():
            # max 3 mode
            if idx<3:
                # mode
                mode = row["Mode"]
                # all stack map
                all_img_path = row["ALL_img_paths"]
                # sample
                lotid_1 = row["Sample1_lotid"]
                img_1 = row["Samp1_img_paths"]
                lotid_2 = row["Sample2_lotid"]
                img_2 = row["Samp2_img_paths"]

                # all image ------------------------
                slide.shapes.add_picture(all_img_path,
                                        Inches(left+0.1),
                                        Inches(top + 1.1*idx),
                                        Inches(width),
                                        Inches(height)
                                        )
                # graph text, mode
                slide = self.add_graph_text(slide=slide,
                                            text=mode,
                                            left=left,
                                            top=top + 1.1*idx - 0.2,
                                            width=width,
                                            height=0.2,
                                            bold=True
                                            )
                # graph text, ALL
                slide = self.add_graph_text(slide=slide,
                                            text="ALL",
                                            left=left+0.1,
                                            top=top + height + 1.1*idx,
                                            width=width,
                                            height=0.2)

                # sample1 image -----------------------
                slide.shapes.add_picture(img_1,
                                        Inches(left+1.3),
                                        Inches(top + 1.1*idx),
                                        Inches(width),
                                        Inches(height)
                                        )
                # graph text, sample1 lot
                slide = self.add_graph_text(slide=slide,
                                            text=lotid_1,
                                            left=left+1.3,
                                            top=top + height + 1.1*idx,
                                            width=width,
                                            height=0.2)

                # sample2 image ------------------------
                slide.shapes.add_picture(img_2,
                                        Inches(left+2.3),
                                        Inches(top + 1.1*idx),
                                        Inches(width),
                                        Inches(height)
                                        )
                # graph text, sample1 lot
                slide = self.add_graph_text(slide=slide,
                                            text=lotid_2,
                                            left=left+2.3,
                                            top=top + height + 1.1*idx,
                                            width=width,
                                            height=0.2)

        # lot samples
        # add text box, left -> top -> width -> height
        textbox3 = slide.shapes.add_textbox(Inches(7),
                                            Inches(1),
                                            Inches(3),
                                            Inches(0.3))
        # add text
        textbox3.text_frame.paragraphs[0].text = "Lot sample graph"
        textbox3.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox3.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox3.text_frame.paragraphs[0].font.bold = True # bold
        textbox3.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # position
        left = 7
        top = 1.5
        width = 3
        height = 1.8
        for idx, row in graph_images_df.iterrows():
            # mode
            mode = row["Mode"]
            # sample
            lotid_1 = row["Sample1_lotid"]
            graph_1 = row["Samp1_graph_paths"]
            lotid_2 = row["Sample2_lotid"]
            graph_2 = row["Samp2_graph_paths"]
            # graph text, mode
            slide = self.add_graph_text(slide=slide,
                                        text=mode,
                                        left=left,
                                        top=top + (height+0.1)*idx,
                                        width=width,
                                        height=0.2,
                                        bold=True
                                        )
            # sample1 image -----------------------
            slide.shapes.add_picture(graph_1,
                                    Inches(left),
                                    Inches(top + (height+0.1)*idx+0.1),
                                    Inches(width),
                                    Inches(height)
                                    )
            # graph text, sample1 lot
            slide = self.add_graph_text(slide=slide,
                                        text=lotid_1,
                                        left=left,
                                        top=top + (height+0.1)*(idx+1)-0.05,
                                        width=width,
                                        height=0.2,
                                        position="center")
            # sample2 image -----------------------
            slide.shapes.add_picture(graph_2,
                                    Inches(left + width),
                                    Inches(top + (height+0.1)*idx + 0.1),
                                    Inches(width),
                                    Inches(height)
                                    )
            # graph text, sample1 lot
            slide = self.add_graph_text(slide=slide,
                                        text=lotid_2,
                                        left=left + width,
                                        top=top + (height+0.1)*(idx+1)-0.05,
                                        width=width,
                                        height=0.2,
                                        position="center")

    ##########################################
    # 4th slide, Auto factor analysis result
    ##########################################
    def importance_analysis_result(self,
                                   test_name:str,
                                   importance_result_df:pd.DataFrame,
                                   importance_graph_images_df:pd.DataFrame):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # ------------------------
        # title
        # ------------------------
        title = slide.shapes.title
        title.text = f"Test : {test_name}, Auto Factor Analysis (Prediction by ML)"
        # title text format and position
        title = self.format_title(title=title)

        # --------------------------
        # sub title by text box
        # --------------------------
        self.add_subtitle(slide=slide,
                          subtitle_sentence="Please input comment.")

        # --------------------------
        # Importance result
        # --------------------------
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(Inches(0.5),
                                            Inches(1),
                                            Inches(3),
                                            Inches(0.3))
        # add text
        textbox1.text_frame.paragraphs[0].text = "Importance rank"
        textbox1.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # importance rank
        # add dataframe
        self.add_matrix_3(slide=slide,
                        df=importance_result_df,
                        left=0.5,
                        top=1.5,
                        width=6,
                        height=15)

        # importance graphs
        # add text box, left -> top -> width -> height
        textbox2 = slide.shapes.add_textbox(Inches(6.8),
                                            Inches(1),
                                            Inches(3),
                                            Inches(0.3))
        # add text
        textbox2.text_frame.paragraphs[0].text = "Top 3 graphs"
        textbox2.text_frame.paragraphs[0].font.size = Pt(12) # font size
        textbox2.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox2.text_frame.paragraphs[0].font.bold = True # bold
        textbox2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # position
        left = 7
        top = 1.5
        width = 3
        height = 1.8
        for idx, row in importance_graph_images_df.iterrows():
            Parameter = row["Parameter"]
            graph_1 = row["graph_1"]
            graph_2 = row["graph_2"]
            # graph text mode
            slide = self.add_graph_text(slide=slide,
                                        text=Parameter,
                                        left=left,
                                        top=top + (height+0.1)*idx,
                                        width=width,
                                        height=0.2,
                                        bold=True
                                        )
            # sample 1 image ---------------------------
            slide.shapes.add_picture(graph_1,
                                    Inches(left),
                                    Inches(top + (height+0.1)*idx+0.1),
                                    Inches(width),
                                    Inches(height)
                                    )
            # graph text, sample1 lot
            slide = self.add_graph_text(slide=slide,
                                        text="Parameter vs target",
                                        left=left,
                                        top=top + (height+0.1)*(idx+1)-0.05,
                                        width=width,
                                        height=0.2,
                                        position="center")
            # sample2 image -----------------------
            slide.shapes.add_picture(graph_2,
                                    Inches(left + width),
                                    Inches(top + (height+0.1)*idx + 0.1),
                                    Inches(width),
                                    Inches(height)
                                    )
            # graph text, sample1 lot
            slide = self.add_graph_text(slide=slide,
                                        text="Trend",
                                        left=left + width,
                                        top=top + (height+0.1)*(idx+1)-0.05,
                                        width=width,
                                        height=0.2,
                                        position="center")

    ##########################################
    # 5th slide, Process image
    ##########################################
    def test_importanct_process(self,
                                test_name:str,
                                process_image:str):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # ------------------------
        # title
        # ------------------------
        title = slide.shapes.title
        title.text = f"Test : {test_name}, Position of suspect process"
        # title text format and position
        title = self.format_title(title=title)

        # --------------------------
        # sub title by text box
        # --------------------------
        self.add_subtitle(slide=slide,
                          subtitle_sentence="Please input comment.")

        # --------------------------
        # graph image
        # --------------------------
        # add slide
        slide.shapes.add_picture(process_image,
                                 Inches(0.5),
                                 Inches(1),
                                 Inches(12.3),
                                 Inches(6.3)
                                 )


    def format_title(self, title):

        # title text format
        title.text_frame.paragraphs[0].font.size = Pt(18) # font size
        title.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        title.text_frame.paragraphs[0].font.bold = True # bold
        title.text_frame.paragraphs[0].font.underline = True # under line
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # title text position, left -> top -> width -> height
        title.left = Inches(0.2)
        title.top = Inches(0.1)
        title.width = Inches(13)
        title.height = Inches(0.5)

        return title

    def add_subtitle(self,
                     slide,
                     subtitle_sentence:str):
        # define position
        left = Inches(0.2)
        top = Inches(0.5)
        width = Inches(13)
        height = Inches(0.5)
        # add text box, left -> top -> width -> height
        textbox = slide.shapes.add_textbox(left, top, width, height)
        # add text
        textbox.text_frame.paragraphs[0].text = subtitle_sentence
        # format
        textbox.text_frame.paragraphs[0].font.size = Pt(18) # font size
        textbox.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        return slide

    # add matrixs, basic with same column width
    def add_matrix_1(self,
                   slide,
                   df:pd.DataFrame,
                   left:float,
                   top:float,
                   width:float,
                   height:float,
                   column_color:tuple,
                   merge_columns:list=None):

        # add table
        table = slide.shapes.add_table(df.shape[0]+1, # height dicrection shape with column
                                       df.shape[1], # width direction shape
                                       Inches(left), # left position
                                       Inches(top), # top position
                                       Inches(width), # width
                                       Inches(height) # height
                                       ).table
        # format the hearder row
        for col_index, col_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = col_name # column names
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(column_color[0],column_color[1],column_color[2]) # Dark blue background
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        # format the data rows
        for row_index, row in df.iterrows():
            for col_index, item in enumerate(row):
                cell = table.cell(row_index+1, col_index) # row is +1 because including column
                cell.text = str(item)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
        # set table height
        row_height = Inches(0.3)
        for row in table.rows:
            row.height = row_height

        # merge columns
        if merge_columns!=None:
            # each data
            for cc in merge_columns:
                col_idx = cc[0]
                col = cc[1]
                # unique value
                unique = df[col].unique()
                for uq in unique:
                    d = df[df[col]==uq][col]
                    if len(d)>0:
                        min_idx = np.min(d.index)+1
                        max_idx = np.max(d.index)+1
                        # execute merge cell
                        merged_cell = table.cell(min_idx, col_idx).merge(table.cell(max_idx, col_idx))
                        table.cell(min_idx, col_idx).text = str(uq)
                        table.cell(min_idx, col_idx).fill.solid()
                        table.cell(min_idx, col_idx).fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.name = 'Meiryo UI'
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        else:
            pass

        return slide

    # add matrixs, yield diagram
    def add_matrix_2(self,
                   slide,
                   df:pd.DataFrame,
                   left:float,
                   top:float,
                   width:float,
                   height:float,
                   merge_columns:list=None):

        # add table
        table = slide.shapes.add_table(df.shape[0]+1, # height dicrection shape with column
                                       df.shape[1], # width direction shape
                                       Inches(left), # left position
                                       Inches(top), # top position
                                       Inches(width), # width
                                       Inches(height) # height
                                       ).table
        # format the hearder row
        for col_index, col_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = col_name # column names
            cell.fill.solid()
            if (col_name=="Total") or (col_name=="Y*") or (col_name=="BIN") or (col_name=="TEST"):
                cell.fill.fore_color.rgb = RGBColor(orange[0],orange[1],orange[2]) # orange back ground
            elif (col_name=="Map_classification") or (col_name=="Slot_dependency"):
                cell.fill.fore_color.rgb = RGBColor(blue[0],blue[1],blue[2]) # blue back ground
            else:
                cell.fill.fore_color.rgb = RGBColor(sky[0],sky[1],sky[2]) # blue back ground
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        # format the data rows
        for row_index, row in df.iterrows():
            for col_index, item in enumerate(row):
                cell = table.cell(row_index+1, col_index) # row is +1 because including column
                cell.text = str(item)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                cell.text_frame.paragraphs[0].font.size = Pt(8)
                cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
        # set table height
        row_height = Inches(0.3)
        for row in table.rows:
            row.height = row_height
        # set table column
        col_width = [0.75, 0.5, 1, 1.5, 2.5, 2.5, 4]
        for idx, col in enumerate(table.columns):
            col.width = Inches(col_width[idx])

        # merge columns
        if merge_columns!=None:
            # each data
            for cc in merge_columns:
                col_idx = cc[0]
                col = cc[1]
                # unique value
                unique = df[col].unique()
                for uq in unique:
                    d = df[df[col]==uq][col]
                    if len(d)>0:
                        min_idx = np.min(d.index)+1
                        max_idx = np.max(d.index)+1
                        # execute merge cell
                        merged_cell = table.cell(min_idx, col_idx).merge(table.cell(max_idx, col_idx))
                        table.cell(min_idx, col_idx).text = str(uq)
                        table.cell(min_idx, col_idx).fill.solid()
                        table.cell(min_idx, col_idx).fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.name = 'Meiryo UI'
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        else:
            pass

        return slide

    # add matrixs, importance with unique column width
    def add_matrix_3(self,
                    slide,
                    df:pd.DataFrame,
                    left:float,
                    top:float,
                    width:float,
                    height:float,
                    merge_columns=None):

        # add table
        table = slide.shapes.add_table(df.shape[0]+1, # height dicrection shape with column
                                       df.shape[1], # width direction shape
                                       Inches(left), # left position
                                       Inches(top), # top position
                                       Inches(width), # width
                                       Inches(height) # height
                                       ).table
        # format the hearder row
        for col_index, col_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = col_name # column names
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(sky[0],sky[1],sky[2]) # blue back ground
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        # format the data rows
        for row_index, row in df.iterrows():
            for col_index, item in enumerate(row):
                cell = table.cell(row_index+1, col_index) # row is +1 because including column
                cell.text = str(item)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                cell.text_frame.paragraphs[0].font.size = Pt(8)
                cell.text_frame.paragraphs[0].font.name = 'Meiryo UI'
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
        # set table height
        row_height = Inches(0.3)
        for row in table.rows:
            row.height = row_height
        # set table column
        col_width = [1, 3, 1, 1]
        for idx, col in enumerate(table.columns):
            col.width = Inches(col_width[idx])

        # merge columns
        if merge_columns!=None:
        # each data
            for cc in merge_columns:
                col_idx = cc[0]
                col = cc[1]
                # unique value
                unique = df[col].unique()
                for uq in unique:
                    d = df[df[col]==uq][col]
                    if len(d)>0:
                        min_idx = np.min(d.index)+1
                        max_idx = np.max(d.index)+1
                        # execute merge cell
                        merged_cell = table.cell(min_idx, col_idx).merge(table.cell(max_idx, col_idx))
                        table.cell(min_idx, col_idx).text = str(uq)
                        table.cell(min_idx, col_idx).fill.solid()
                        table.cell(min_idx, col_idx).fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.name = 'Meiryo UI'
                        table.cell(min_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        else:
            pass

        return slide

    # add graph text
    def add_graph_text(self,
                       slide,
                       text:str,
                       left:float,
                       top:float,
                       width:float,
                       height:float,
                       bold:bool=False,
                       position:str="left"):

        # add text box
        textbox = slide.shapes.add_textbox(Inches(left),
                                           Inches(top),
                                           Inches(width),
                                           Inches(height))
        # add text
        textbox.text_frame.paragraphs[0].text = text
        textbox.text_frame.paragraphs[0].font.size = Pt(7) # font size
        textbox.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        # bold
        if bold:
            textbox.text_frame.paragraphs[0].font.bold = True # bold
        # text position
        if position=="left":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position
        elif position=="right":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        elif position=="center":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def save_pptx(self):
        self.prs.save(self.savefilename)

# for debug
if __name__=="__main__":

    # instance
    mp = MakePPT(savefilename="report.pptx")

    # add 1st slide
    mp.data_summary(dataset_summary=sd.dataset,
                    ml_learning_dataset_summary=sd.mllearning)
    # add 2nd slide
    mp.yield_diagram(diagram=sd.diagram)

    for i in range(10):
        # report parameters
        # (1) test name
        # (2) test value stats matrix
        # (3) test value distribution graph path
        # (4) condition analysis result dataframe
        # (5) map image path dataset dataframe
        # (6) lot sample dataset dataframe
        # (7) importance result dataframe
        # (8) importance graph sample dataframe
        # (9) process flow analysis image path
        # add 3rd slide
        mp.test_data_condition(test_name=f"TEST{i}",
                                stats_matrix=sd.stat_matrix,
                                stats_histgram=sd.graph,
                                condition_analysis=sd.cond,
                                map_images_df=sd.map_images,
                                graph_images_df=sd.graph_lot_images)
        # add 4th slide
        mp.importance_analysis_result(test_name=f"TEST{i}",
                                        importance_result_df=sd.imp,
                                        importance_graph_images_df=sd.imp_graphs)
        # add 5th slide
        mp.test_importanct_process(test_name=f"TEST{i}",
                                    process_image=sd.process_flow)
        # save file
    mp.save_pptx()