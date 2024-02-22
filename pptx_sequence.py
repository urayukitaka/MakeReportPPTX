# Libraries
# pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
# pandas
import pandas as pd
import numpy as np
# datetime
from datetime import datetime as dt


class PPT_SEQ:

    def __init__(self, savefilename:str):

        # init
        self.savefilename = savefilename
        # new presentation
        self.prs = Presentation()
        # wide screen
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    # text
    def make_sequence(self,
                      seq_df:pd.DataFrame,
                      category:list
                      ):

        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "Process sequence analysis"
        # title text format and position
        title = self.format_title(title=title)

        # start x, y
        start_x, start_y = Inches(0.5), Inches(0.7)
        text_box_h = Inches(0.2)
        box_start_x, box_start_y = Inches(0.7), start_y
        box_w, box_h = Inches(0.6), Inches(0.2)

        # ------------------------
        # Node contents
        # ------------------------
        x_counter = 0
        y_counter = 0
        for cat in category:
            # Category name
            # add text box, left -> top -> width -> height
            textbox1 = slide.shapes.add_textbox(start_x + x_counter*Inches(1.0),
                                                start_y + y_counter*Inches(0.21),
                                                Inches(1.5),
                                                text_box_h)
            # add text
            textbox1.text_frame.paragraphs[0].text = f"{cat}"
            textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
            textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
            textbox1.text_frame.paragraphs[0].font.bold = True # bold
            textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position
            # y counter
            y_counter += 1

            # slice dataset
            sdf = seq_df[seq_df["M_Process"]==cat].reset_index(drop=True)
            # each data
            for idx, row in sdf.iterrows():
                # y counter
                if (idx!=0)&((idx+1)%12==1):
                    y_counter += 1
                # process name
                process = row["Process"]
                flg = row["flg"]
                # dimension
                x, y = box_start_x + (x_counter%12)*Inches(1.0), box_start_y + y_counter*Inches(0.21)
                cx, cy = box_w, box_h
                end_x1, end_y1 = x + box_w, y
                # try add box
                try:
                    parent_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x, y, cx, cy
                    )
                    parent_shape.text = "{}".format(process)
                    parent_shape.line.fill.background()
                    fill = parent_shape.fill
                    fill.solid()
                    if flg==1:
                        fill.fore_color.rgb = RGBColor(255, 0, 0)
                        parent_text_frame = parent_shape.text_frame
                        parent_text_frame.paragraphs[0].font.size = Pt(7)
                        parent_text_frame.paragraphs[0].font.bold = True
                        parent_text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                    else:
                        fill.fore_color.rgb = RGBColor(255, 255, 255)
                        parent_text_frame = parent_shape.text_frame
                        parent_text_frame.paragraphs[0].font.size = Pt(7)
                        parent_text_frame.paragraphs[0].font.bold = False
                        parent_text_frame.paragraphs[0].font.color.rgb = RGBColor(32,32,32)

                    # add arrow
                    if idx+1 != len(sdf):
                        arw = slide.shapes.add_connector(
                            MSO_CONNECTOR.STRAIGHT,
                            end_x1,
                            end_y1 + int(box_h/2),
                            end_x1 + Inches(0.4),
                            end_y1 + int(box_h/2)
                        )
                        arw.line.width = Pt(0.75)
                        arw.line.color.rgb = RGBColor(32, 32, 32)
                except:
                    pass
                # add x counter
                x_counter += 1

            # x counter, reset
            x_counter = 0
            # y counter
            y_counter += 1

    # basic title
    def format_title(self, title):

        # title text format
        title.text_frame.paragraphs[0].font.size = Pt(18) # font size
        title.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        title.text_frame.paragraphs[0].font.bold = True # bold
        title.text_frame.paragraphs[0].font.underline = True # under line
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # title text position, left -> top -> width -> height
        title.left = Inches(0.2)
        title.top = Inches(0.1)
        title.width = Inches(13)
        title.height = Inches(0.5)

        return title

    def add_subtitle(self,
                     slide,
                     subtitle_sentence:str):
        # define position
        left = Inches(0.2)
        top = Inches(0.5)
        width = Inches(13)
        height = Inches(0.5)
        # add text box, left -> top -> width -> height
        textbox = slide.shapes.add_textbox(left, top, width, height)
        # add text
        textbox.text_frame.paragraphs[0].text = subtitle_sentence
        # format
        textbox.text_frame.paragraphs[0].font.size = Pt(18) # font size
        textbox.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        return slide

    def save_pptx(self):
        self.prs.save(self.savefilename)
