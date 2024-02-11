# Libraries
# pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
# pandas
import pandas as pd
import numpy as np
# datetime
from datetime import datetime as dt

# dataset sample
import sample_dataset as sd

# parameter
orange = (165,42,42)
blue = (0,0,128)
sky = (0,153,255)
comment = "注意事項\n・\n・"

mapimg = r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\map.png"

# class
class PPT_DIM:

    def __init__(self, savefilename:str):

        # init
        self.savefilename = savefilename
        # new presentation
        self.prs = Presentation()
        # wide screen
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    # text
    def make_diagram(self,
                     diagram:pd.DataFrame,
                     color_map:dict):

        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "Yield diagram"
        # title text format and position
        title = self.format_title(title=title)

        # ------------------------
        # Node contents
        # ------------------------
        c_start_x, c_start_y = Inches(2.3), Inches(0.5)
        text_box_h = Inches(0.3)
        # BIN name
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(c_start_x,
                                            c_start_y,
                                            Inches(1.2),
                                            text_box_h)
        # add text
        textbox1.text_frame.paragraphs[0].text = "BIN Failure rate"
        textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # TEST name
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(c_start_x + Inches(1.5),
                                            c_start_y,
                                            Inches(2.2),
                                            text_box_h)
        # add text
        textbox1.text_frame.paragraphs[0].text = "TEST Failure rate"
        textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # MAP mode
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(c_start_x + Inches(3.75),
                                            c_start_y,
                                            Inches(1.25),
                                            text_box_h)
        # add text
        textbox1.text_frame.paragraphs[0].text = "MAP and mode"
        textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # SLOT trend
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(c_start_x + Inches(5.15),
                                            c_start_y,
                                            Inches(1),
                                            text_box_h)
        # add text
        textbox1.text_frame.paragraphs[0].text = "Slot trend"
        textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # Factor Analysis result
        # add text box, left -> top -> width -> height
        textbox1 = slide.shapes.add_textbox(c_start_x + Inches(6.2),
                                            c_start_y,
                                            Inches(4.6),
                                            text_box_h)
        # add text
        textbox1.text_frame.paragraphs[0].text = "Factor analysis result"
        textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
        textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox1.text_frame.paragraphs[0].font.bold = True # bold
        textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # dim
        start_left, start_top = Inches(0.1), Inches(0.8)
        box_h = Inches(0.3)

        # ------------------------
        # 1st, yield
        # ------------------------
        yld = diagram["Yield"].values[0]
        # node
        x, y, cx, cy = start_left, start_top, Inches(0.8), int(box_h)
        # end point
        end_x1, end_y1 = start_left + Inches(0.8), start_top
        try:
            parent_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, cx, cy
            )
            parent_shape.text = "Yield {}%".format(yld)
            parent_shape.line.fill.background()
            fill = parent_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            parent_text_frame = parent_shape.text_frame
            parent_text_frame.paragraphs[0].font.size = Pt(10)
            parent_text_frame.paragraphs[0].font.bold = True
            parent_text_frame.paragraphs[0].font.color.rgb = RGBColor(32,32,32)
        except:
            pass

        # ------------------------
        # 2nd, yield type
        # ------------------------
        cntr2 = 0
        for idx2, y_ in enumerate(diagram["Y*"].drop_duplicates().values):
            # Bin failure rate
            yld_ = diagram[(diagram["Y*"]==y_)]["Y*_FailureRate"].values[0]
            # node
            x, y, cx, cy = end_x1 + Inches(0.3), end_y1 + cntr2*(box_h + Inches(0.05)), Inches(0.8), int(box_h)
            # end point
            end_x2, end_y2 = x + Inches(0.8), y
            try:
                node2_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y, cx, cy
                )
                node2_shape.text = "{}_{}%".format(y_, yld_)
                node2_shape.line.fill.background()
                fill = node2_shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                node2_text_frame = node2_shape.text_frame
                node2_text_frame.paragraphs[0].font.size = Pt(10)
                node2_text_frame.paragraphs[0].font.bold = True
                node2_text_frame.paragraphs[0].font.color.rgb = RGBColor(32,32,32)
                # add line
                line12 = slide.shapes.add_connector(
                    MSO_CONNECTOR.ELBOW,
                    end_x1,
                    end_y1+int(box_h/2),
                    x,
                    y+int(box_h/2)
                )
                line12.line.width = Pt(0.75)
                line12.line.color.rgb = RGBColor(32, 32, 32)
            except:
                pass

            # ------------------------
            # 3rd, BIN failure rate
            # ------------------------
            cntr3 = 0
            for idx3, b in enumerate(diagram[(diagram["Y*"]==y_)]["BIN"].drop_duplicates().values):
                # Bin failure rate
                bfr = diagram[(diagram["Y*"]==y_)&(diagram["BIN"]==b)]["BIN_FailureRate"].values[0]
                # node
                x, y, cx, cy = end_x2 + Inches(0.3), end_y2+ cntr3*(box_h + Inches(0.05)), Inches(1.2), int(box_h)
                # end point
                end_x3, end_y3 = x + Inches(1.2), y
                try:
                    node3_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x, y, cx, cy
                    )
                    node3_shape.text = "{}_{}%".format(b, bfr)
                    node3_shape.line.fill.background()
                    fill = node3_shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(color_map[b][0], color_map[b][1], color_map[b][2])
                    node3_text_frame = node3_shape.text_frame
                    node3_text_frame.paragraphs[0].font.size = Pt(10)
                    # add line
                    line23 = slide.shapes.add_connector(
                        MSO_CONNECTOR.ELBOW,
                        end_x2,
                        end_y2+int(box_h/2),
                        x,
                        y+int(box_h/2)
                    )
                    line23.line.width = Pt(0.75)
                    line23.line.color.rgb = RGBColor(32, 32, 32)
                except:
                    pass

                # ------------------------
                # 4th, TEST failure rate
                # ------------------------
                for idx4, t in enumerate(diagram[(diagram["Y*"]==y_)&(diagram["BIN"]==b)]["TEST"].drop_duplicates().values):
                    # sampling
                    sample = diagram[(diagram["Y*"]==y_)&(diagram["BIN"]==b)&(diagram["TEST"]==t)]
                    tfr = sample["TEST_FailureRate"].values[0]
                    mimg = sample["MAP_img"].values[0]
                    mclf = sample["MAP_classification"].values[0]
                    strd = sample["SlotTrend"].values[0]
                    imp = sample["ImportanceAnalysis"].values[0]
                    # ------
                    # node
                    # ------
                    x, y, cx, cy = end_x3 + Inches(0.3), end_y3 + idx4*(box_h + Inches(0.05)), Inches(2.2), int(box_h)
                    # end point
                    end_x4 = x + Inches(2.2)
                    try:
                        node4_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            x, y, cx, cy
                        )
                        node4_shape.text = "{}_{}%".format(t, tfr)
                        node4_shape.line.fill.background()
                        fill = node4_shape.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(color_map[b][0], color_map[b][1], color_map[b][2])
                        node4_text_frame = node4_shape.text_frame
                        node4_text_frame.paragraphs[0].font.size = Pt(8)
                        # add line
                        line34 = slide.shapes.add_connector(
                            MSO_CONNECTOR.ELBOW,
                            end_x3,
                            end_y3+int(box_h/2),
                            x,
                            y+int(box_h/2)
                        )
                        line34.line.width = Pt(0.75)
                        line34.line.color.rgb = RGBColor(32, 32, 32)
                    except:
                        pass

                    # ------
                    # map image
                    # ------
                    slide.shapes.add_picture(mimg,
                                    end_x4 + Inches(0.05),
                                    y,
                                    box_h,
                                    box_h
                                    )
                    # end point
                    end_x5 = end_x4 + Inches(0.05) + box_h

                    # ------
                    # map clf
                    # ------
                    x, y, cx, cy = end_x5 + Inches(0.05), y, Inches(1.), int(box_h)
                    # end point
                    end_x6 = x + Inches(1.)
                    try:
                        if not pd.isna(mclf):
                            node6_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                x, y, cx, cy
                            )
                            node6_shape.line.fill.background()
                            node6_shape.text = "{}".format(mclf)
                            fill = node6_shape.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255)
                            node6_text_frame = node6_shape.text_frame
                            node6_text_frame.paragraphs[0].font.size = Pt(8)
                            node6_text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                    except:
                        pass

                    # ------
                    # slot
                    # ------
                    x, y, cx, cy = end_x6 + Inches(0.05), y, Inches(1.), int(box_h)
                    # end point
                    end_x7 = x + Inches(1.)
                    try:
                        if not pd.isna(strd):
                            node7_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                x, y, cx, cy
                            )
                            node7_shape.text = "{}".format(strd)
                            node7_shape.line.fill.background()
                            fill = node7_shape.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255)
                            node7_text_frame = node7_shape.text_frame
                            node7_text_frame.paragraphs[0].font.size = Pt(8)
                            node7_text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                    except:
                        pass

                    # ------
                    # Importance
                    # ------
                    x, y, cx, cy = end_x7 + Inches(0.05), y, Inches(4.6), int(box_h)
                    try:
                        if not pd.isna(imp):
                            node8_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                x, y, cx, cy
                            )
                            node8_shape.text = "{}".format(imp)
                            node8_shape.line.fill.background()
                            fill = node8_shape.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255)
                            node8_text_frame = node8_shape.text_frame
                            node8_text_frame.paragraphs[0].font.size = Pt(8)
                            node8_text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                    except:
                        pass

                    # counter
                    cntr2 += 1
                    cntr3 += 1

    def bin_diagram(self,
                    binname:str,
                    diagram:pd.DataFrame,
                    color_map:dict):

        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "BIN {} diagram"
        # title text format and position
        title = self.format_title(title=title)

    # image
    def add_img(self
                ):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "Image dimension check"
        # title text format and position
        title = self.format_title(title=title)

        # sub title
        self.add_subtitle(slide=slide,
                          subtitle_sentence = "Sub title")

        # dimensions
        dims = []
        for i in range(0,14,1):
            for j in range(0,9,1):
                dims.append((i,j))

        for d in dims:
            # image
            slide.shapes.add_picture(mapimg,
                                     Inches(d[0]),
                                     Inches(d[1]),
                                     Inches(0.5),
                                     Inches(0.5)
                                     )

            # add text box, left -> top -> width -> height
            textbox1 = slide.shapes.add_textbox(Inches(d[0]),
                                                Inches(d[1]+0.5),
                                                Inches(1),
                                                Inches(0.2))
            # add text
            textbox1.text_frame.paragraphs[0].text = "({},{})".format(d[0],d[1])
            textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
            textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
            textbox1.text_frame.paragraphs[0].font.bold = False # bold
            textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position


    # basic title
    def format_title(self, title):

        # title text format
        title.text_frame.paragraphs[0].font.size = Pt(18) # font size
        title.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        title.text_frame.paragraphs[0].font.bold = True # bold
        title.text_frame.paragraphs[0].font.underline = True # under line
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # title text position, left -> top -> width -> height
        title.left = Inches(0.2)
        title.top = Inches(0.1)
        title.width = Inches(13)
        title.height = Inches(0.5)

        return title

    def add_subtitle(self,
                     slide,
                     subtitle_sentence:str):
        # define position
        left = Inches(0.2)
        top = Inches(0.5)
        width = Inches(13)
        height = Inches(0.5)
        # add text box, left -> top -> width -> height
        textbox = slide.shapes.add_textbox(left, top, width, height)
        # add text
        textbox.text_frame.paragraphs[0].text = subtitle_sentence
        # format
        textbox.text_frame.paragraphs[0].font.size = Pt(18) # font size
        textbox.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        return slide

    def save_pptx(self):
        self.prs.save(self.savefilename)

# for debug
if __name__ == "__main__":

    # instance
    mp = PPT_DIM(
        savefilename = "dim_test.pptx"
    )

    # text
    mp.add_text()

    # image
    mp.add_img()

    # save
    mp.save_pptx()