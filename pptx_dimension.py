# Libraries
# pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# pandas
import pandas as pd
import numpy as np
# datetime
from datetime import datetime as dt

# dataset sample
import sample_dataset as sd

# parameter
orange = (165,42,42)
blue = (0,0,128)
sky = (0,153,255)
comment = "注意事項\n・\n・"

mapimg = r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\map.png"

# class
class PPT_DIM:

    def __init__(self, savefilename:str):

        # init
        self.savefilename = savefilename
        # new presentation
        self.prs = Presentation()
        # wide screen
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    # text
    def add_text(self,
                 ):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "Text dimension check"
        # title text format and position
        title = self.format_title(title=title)

        # sub title
        self.add_subtitle(slide=slide,
                          subtitle_sentence = "Sub title")

        # dimensions
        dims = []
        for i in range(0,14,1):
            for j in range(0,9,1):
                dims.append((i,j))

        for d in dims:
            # add text box, left -> top -> width -> height
            textbox1 = slide.shapes.add_textbox(Inches(d[0]),
                                                Inches(d[1]),
                                                Inches(1),
                                                Inches(0.2))
            # add text
            textbox1.text_frame.paragraphs[0].text = "({},{})".format(d[0],d[1])
            textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
            textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
            textbox1.text_frame.paragraphs[0].font.bold = False # bold
            textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position


    # image
    def add_img(self
                ):
        # select slide layout
        slide_layout = self.prs.slide_layouts[5] # 5 is blank slide
        # add slide
        slide = self.prs.slides.add_slide(slide_layout)

        # title
        title = slide.shapes.title
        title.text = "Image dimension check"
        # title text format and position
        title = self.format_title(title=title)

        # sub title
        self.add_subtitle(slide=slide,
                          subtitle_sentence = "Sub title")

        # dimensions
        dims = []
        for i in range(0,14,1):
            for j in range(0,9,1):
                dims.append((i,j))

        for d in dims:
            # image
            slide.shapes.add_picture(mapimg,
                                     Inches(d[0]),
                                     Inches(d[1]),
                                     Inches(0.5),
                                     Inches(0.5)
                                     )

            # add text box, left -> top -> width -> height
            textbox1 = slide.shapes.add_textbox(Inches(d[0]),
                                                Inches(d[1]+0.5),
                                                Inches(1),
                                                Inches(0.2))
            # add text
            textbox1.text_frame.paragraphs[0].text = "({},{})".format(d[0],d[1])
            textbox1.text_frame.paragraphs[0].font.size = Pt(8) # font size
            textbox1.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
            textbox1.text_frame.paragraphs[0].font.bold = False # bold
            textbox1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position


    # basic title
    def format_title(self, title):

        # title text format
        title.text_frame.paragraphs[0].font.size = Pt(18) # font size
        title.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        title.text_frame.paragraphs[0].font.bold = True # bold
        title.text_frame.paragraphs[0].font.underline = True # under line
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        # title text position, left -> top -> width -> height
        title.left = Inches(0.2)
        title.top = Inches(0.1)
        title.width = Inches(13)
        title.height = Inches(0.5)

        return title

    def add_subtitle(self,
                     slide,
                     subtitle_sentence:str):
        # define position
        left = Inches(0.2)
        top = Inches(0.5)
        width = Inches(13)
        height = Inches(0.5)
        # add text box, left -> top -> width -> height
        textbox = slide.shapes.add_textbox(left, top, width, height)
        # add text
        textbox.text_frame.paragraphs[0].text = subtitle_sentence
        # format
        textbox.text_frame.paragraphs[0].font.size = Pt(18) # font size
        textbox.text_frame.paragraphs[0].font.name = "Meiryo UI" # font
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # position

        return slide

    def save_pptx(self):
        self.prs.save(self.savefilename)

# for debug
if __name__ == "__main__":

    # instance
    mp = PPT_DIM(
        savefilename = "dim_test.pptx"
    )

    # text
    mp.add_text()

    # image
    mp.add_img()

    # save
    mp.save_pptx()